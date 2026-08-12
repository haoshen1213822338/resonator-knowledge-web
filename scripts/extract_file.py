"""Extract readable text from common office documents.

The web app calls this script after a file upload. It converts supported
binary formats into plain Markdown-like text so the AI organizer can work from
local extracted content instead of raw files.
"""

from __future__ import annotations

import argparse
import json
import os
import shutil
import subprocess
import sys
import tempfile
import textwrap
import zipfile
from pathlib import Path
from typing import Any


SUPPORTED_EXTENSIONS = {
    ".docx",
    ".pdf",
    ".xlsx",
    ".pptx",
    ".md",
    ".txt",
    ".csv",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".bmp",
    ".mp4",
    ".mov",
    ".mkv",
    ".avi",
    ".webm",
    ".m4v",
}

VIDEO_EXTENSIONS = {".mp4", ".mov", ".mkv", ".avi", ".webm", ".m4v"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp"}

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
if hasattr(sys.stderr, "reconfigure"):
    sys.stderr.reconfigure(encoding="utf-8")


def emit_progress(phase: str, progress: int | float | None = None, detail: str = "") -> None:
    """Write machine-readable progress events to stderr for the Node backend."""

    payload: dict[str, Any] = {
        "type": "progress",
        "phase": phase,
    }
    if progress is not None:
        payload["progress"] = max(0, min(100, int(progress)))
    if detail:
        payload["detail"] = detail
    print(json.dumps(payload, ensure_ascii=False), file=sys.stderr, flush=True)


def read_text_file(path: Path) -> str:
    """Read a text-like file with practical encoding fallbacks."""

    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def extract_docx(path: Path) -> str:
    """Extract paragraphs and tables from a Word document."""

    emit_progress("读取 Word 文档", 18, path.name)
    from docx import Document

    document = Document(path)
    parts: list[str] = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table_index, table in enumerate(document.tables, 1):
        rows: list[str] = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"## 表格 {table_index}\n" + "\n".join(rows))

    return "\n\n".join(parts)


def extract_pdf(path: Path) -> str:
    """Extract text and simple tables from a PDF."""

    emit_progress("读取 PDF 页面", 18, path.name)
    try:
        import pdfplumber
    except ImportError:
        pdfplumber = None

    if pdfplumber is not None:
        parts: list[str] = []
        with pdfplumber.open(path) as pdf:
            for page_index, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    parts.append(f"## 第 {page_index} 页\n{page_text.strip()}")
        return "\n\n".join(parts)

    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts = []
    for page_index, page in enumerate(reader.pages, 1):
        page_text = page.extract_text() or ""
        if page_text.strip():
            parts.append(f"## 第 {page_index} 页\n{page_text.strip()}")
    return "\n\n".join(parts)


def extract_xlsx(path: Path) -> str:
    """Extract visible cell values from an Excel workbook."""

    emit_progress("读取 Excel 工作表", 18, path.name)
    from openpyxl import load_workbook

    workbook = load_workbook(path, data_only=True, read_only=True)
    parts: list[str] = []

    for sheet in workbook.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value).strip() for value in row]
            while values and not values[-1]:
                values.pop()
            if any(values):
                rows.append(" | ".join(values))
        if rows:
            parts.append(f"## 工作表：{sheet.title}\n" + "\n".join(rows))

    return "\n\n".join(parts)


def extract_pptx(path: Path) -> str:
    """Extract slide text and table cells from a PowerPoint file."""

    emit_progress("读取 PPT 页面", 18, path.name)
    from pptx import Presentation

    presentation = Presentation(path)
    parts: list[str] = []

    for slide_index, slide in enumerate(presentation.slides, 1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_parts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    if any(cells):
                        slide_parts.append(" | ".join(cells))
        if slide_parts:
            parts.append(f"## 幻灯片 {slide_index}\n" + "\n\n".join(slide_parts))

    return "\n\n".join(parts)


def save_analysis_image(source: Path, target: Path, max_side: int = 1800) -> Path:
    """Normalize a visual asset for reliable and economical model input."""

    from PIL import Image

    target.parent.mkdir(parents=True, exist_ok=True)
    with Image.open(source) as image:
        image = image.convert("RGB")
        image.thumbnail((max_side, max_side))
        image.save(target, "JPEG", quality=84, optimize=True)
    return target


def render_pdf_pages(path: Path, assets_dir: Path) -> list[dict[str, str]]:
    """Render every PDF page so scanned text, layouts and charts remain visible."""

    try:
        import pypdfium2 as pdfium
    except ImportError as exc:
        raise RuntimeError("扫描 PDF 视觉解析需要安装 pypdfium2。") from exc

    pdf = pdfium.PdfDocument(str(path))
    assets: list[dict[str, str]] = []
    total = max(1, len(pdf))
    for page_index in range(len(pdf)):
        emit_progress(
            "渲染 PDF 页面",
            45 + ((page_index + 1) / total) * 25,
            f"{page_index + 1}/{total}",
        )
        page = pdf[page_index]
        image = page.render(scale=1.6).to_pil().convert("RGB")
        image.thumbnail((1800, 1800))
        target = assets_dir / f"pdf_page_{page_index + 1:04d}.jpg"
        image.save(target, "JPEG", quality=84, optimize=True)
        assets.append({
            "path": str(target),
            "kind": "pdf_page",
            "label": f"第 {page_index + 1} 页",
        })
    return assets


def render_ppt_slides(path: Path, assets_dir: Path) -> list[dict[str, str]]:
    """Render PPT slides through PowerPoint when available, with image fallback."""

    export_dir = assets_dir / "ppt_slides"
    export_dir.mkdir(parents=True, exist_ok=True)
    powershell = shutil.which("powershell") or shutil.which("powershell.exe")
    if powershell and os.name == "nt":
        script_path = assets_dir / "export_ppt.ps1"
        script_path.write_text(
            textwrap.dedent(
                """
                param([string]$InputPath, [string]$OutputDir)
                $app = $null
                $presentation = $null
                try {
                  $app = New-Object -ComObject PowerPoint.Application
                  $presentation = $app.Presentations.Open($InputPath, $true, $true, $false)
                  $presentation.Export($OutputDir, "JPG", 1600, 900)
                } finally {
                  if ($presentation) { $presentation.Close() }
                  if ($app) { $app.Quit() }
                }
                """
            ).strip(),
            encoding="utf-8-sig",
        )
        try:
            run_command([
                powershell,
                "-NoProfile",
                "-ExecutionPolicy",
                "Bypass",
                "-File",
                str(script_path),
                "-InputPath",
                str(path.resolve()),
                "-OutputDir",
                str(export_dir.resolve()),
            ])
        except Exception:
            pass

    slide_images = sorted(
        [*export_dir.glob("*.jpg"), *export_dir.glob("*.jpeg"), *export_dir.glob("*.png")],
        key=lambda item: int("".join(filter(str.isdigit, item.stem)) or 0),
    )
    if slide_images:
        return [
            {"path": str(image), "kind": "ppt_slide", "label": f"幻灯片 {index}"}
            for index, image in enumerate(slide_images, 1)
        ]

    # Linux deployments without LibreOffice/PowerPoint still retain embedded visuals.
    from pptx import Presentation

    presentation = Presentation(path)
    assets: list[dict[str, str]] = []
    for slide_index, slide in enumerate(presentation.slides, 1):
        for shape_index, shape in enumerate(slide.shapes, 1):
            image = getattr(shape, "image", None)
            if image is None:
                continue
            suffix = f".{image.ext.lower()}"
            source = assets_dir / f"slide_{slide_index:04d}_image_{shape_index:03d}{suffix}"
            source.write_bytes(image.blob)
            target = assets_dir / f"slide_{slide_index:04d}_image_{shape_index:03d}.jpg"
            try:
                save_analysis_image(source, target)
            except Exception:
                continue
            assets.append({
                "path": str(target),
                "kind": "ppt_image",
                "label": f"幻灯片 {slide_index} 内嵌图片 {shape_index}",
            })
    return assets


def extract_docx_visuals(path: Path, assets_dir: Path) -> list[dict[str, str]]:
    """Extract embedded Word images for semantic understanding."""

    assets: list[dict[str, str]] = []
    with zipfile.ZipFile(path) as archive:
        media_names = sorted(
            name for name in archive.namelist() if name.startswith("word/media/")
        )
        for index, media_name in enumerate(media_names, 1):
            suffix = Path(media_name).suffix.lower() or ".bin"
            source = assets_dir / f"docx_image_{index:04d}{suffix}"
            source.write_bytes(archive.read(media_name))
            target = assets_dir / f"docx_image_{index:04d}.jpg"
            try:
                save_analysis_image(source, target)
            except Exception:
                continue
            assets.append({
                "path": str(target),
                "kind": "docx_image",
                "label": f"Word 内嵌图片 {index}",
            })
    return assets


def render_xlsx_pages(path: Path, assets_dir: Path) -> list[dict[str, str]]:
    """Render Excel print pages through Excel so charts can be understood visually."""

    powershell = shutil.which("powershell") or shutil.which("powershell.exe")
    if not powershell or os.name != "nt":
        return []

    pdf_path = assets_dir / "excel_render.pdf"
    script_path = assets_dir / "export_excel.ps1"
    script_path.write_text(
        textwrap.dedent(
            """
            param([string]$InputPath, [string]$OutputPath)
            $app = $null
            $workbook = $null
            try {
              $app = New-Object -ComObject Excel.Application
              $app.Visible = $false
              $app.DisplayAlerts = $false
              $workbook = $app.Workbooks.Open($InputPath, 0, $true)
              $workbook.ExportAsFixedFormat(0, $OutputPath)
            } finally {
              if ($workbook) { $workbook.Close($false) }
              if ($app) { $app.Quit() }
            }
            """
        ).strip(),
        encoding="utf-8-sig",
    )
    try:
        run_command([
            powershell,
            "-NoProfile",
            "-ExecutionPolicy",
            "Bypass",
            "-File",
            str(script_path),
            "-InputPath",
            str(path.resolve()),
            "-OutputPath",
            str(pdf_path.resolve()),
        ])
        if not pdf_path.exists():
            return []
        assets = render_pdf_pages(pdf_path, assets_dir / "excel_pages")
        for index, asset in enumerate(assets, 1):
            asset["kind"] = "xlsx_page"
            asset["label"] = f"Excel 打印页面 {index}"
        return assets
    except Exception:
        return []


def extract_image(path: Path, report_progress: bool = True) -> str:
    """Extract visible text from an image by OCR."""

    if report_progress:
        emit_progress("图片 OCR 识别", 22, path.name)
    try:
        from rapidocr_onnxruntime import RapidOCR
    except ImportError as exc:
        raise RuntimeError(
            "当前 Python 环境未安装 OCR 依赖，请先安装 rapidocr-onnxruntime。"
        ) from exc

    engine = RapidOCR()
    result, _ = engine(str(path))
    if not result:
        return ""

    lines: list[str] = []
    for item in result:
        if len(item) < 2:
            continue
        text = str(item[1]).strip()
        score = item[2] if len(item) > 2 else None
        if not text:
            continue
        if isinstance(score, (int, float)):
            lines.append(f"- {text}（置信度：{score:.2f}）")
        else:
            lines.append(f"- {text}")

    return "\n".join(lines)


def require_ffmpeg() -> str:
    """Return the ffmpeg executable path or raise a clear setup error."""

    executable = os.environ.get("FFMPEG_PATH") or shutil.which("ffmpeg")
    if not executable:
        raise RuntimeError(
            "视频解析需要先安装 ffmpeg，并确保命令行可以直接运行 ffmpeg。"
        )
    return executable


def run_command(command: list[str]) -> subprocess.CompletedProcess[str]:
    """Run a command and return text output for diagnostics."""

    try:
        return subprocess.run(
            command,
            check=True,
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
        )
    except subprocess.CalledProcessError as exc:
        detail = (exc.stderr or exc.stdout or str(exc)).strip()
        raise RuntimeError(f"视频处理命令执行失败：{detail[-2000:]}") from exc


def transcribe_audio(audio_path: Path) -> str:
    """Transcribe audio using faster-whisper first, then whisper if installed."""

    model_name = os.environ.get("VIDEO_WHISPER_MODEL", "base")

    try:
      from faster_whisper import WhisperModel
    except ImportError:
      WhisperModel = None

    if WhisperModel is not None:
        model = WhisperModel(model_name, device=os.environ.get("VIDEO_WHISPER_DEVICE", "cpu"))
        segments, _ = model.transcribe(str(audio_path), language="zh", vad_filter=True)
        lines = []
        for segment in segments:
            text = segment.text.strip()
            if text:
                lines.append(f"- {format_timestamp(segment.start)} - {format_timestamp(segment.end)}：{text}")
        return "\n".join(lines)

    try:
        import whisper
    except ImportError as exc:
        raise RuntimeError(
            "视频语音转文字需要安装 faster-whisper 或 openai-whisper。"
        ) from exc

    model = whisper.load_model(model_name)
    result = model.transcribe(str(audio_path), language="zh")
    lines = []
    for segment in result.get("segments", []):
        text = str(segment.get("text", "")).strip()
        if text:
            lines.append(
                f"- {format_timestamp(float(segment.get('start', 0)))} - "
                f"{format_timestamp(float(segment.get('end', 0)))}：{text}"
            )
    return "\n".join(lines)


def format_timestamp(seconds: float) -> str:
    """Format seconds as HH:MM:SS."""

    total = max(0, int(seconds))
    hours, remainder = divmod(total, 3600)
    minutes, secs = divmod(remainder, 60)
    return f"{hours:02d}:{minutes:02d}:{secs:02d}"


def extract_video(path: Path, assets_dir: Path) -> tuple[str, list[dict[str, str]]]:
    """Extract speech transcript and frame OCR from a video."""

    emit_progress("准备视频解析", 5, path.name)
    assets_dir.mkdir(parents=True, exist_ok=True)
    ffmpeg = require_ffmpeg()
    interval = max(1, int(os.environ.get("VIDEO_FRAME_INTERVAL_SECONDS", "10")))
    max_frames = max(1, int(os.environ.get("VIDEO_MAX_KEYFRAMES", "120")))

    with tempfile.TemporaryDirectory(prefix="kb_video_") as temp_dir:
        temp_path = Path(temp_dir)
        audio_path = temp_path / "audio.wav"
        frame_pattern = str(assets_dir / "video_frame_%05d.jpg")

        emit_progress("抽取视频音频", 14, path.name)
        run_command([
            ffmpeg,
            "-y",
            "-i",
            str(path),
            "-vn",
            "-ac",
            "1",
            "-ar",
            "16000",
            str(audio_path),
        ])

        emit_progress("语音转文字", 36, path.name)
        transcript = transcribe_audio(audio_path)

        frame_ocr_lines: list[str] = []
        try:
            emit_progress("抽取关键帧", 68, path.name)
            run_command([
                ffmpeg,
                "-y",
                "-i",
                str(path),
                "-vf",
                f"fps=1/{interval}",
                "-frames:v",
                str(max_frames),
                frame_pattern,
            ])
            frames = sorted(assets_dir.glob("video_frame_*.jpg"))
            total_frames = max(1, len(frames))
            for index, frame in enumerate(frames):
                emit_progress(
                    "关键帧 OCR",
                    74 + (index / total_frames) * 20,
                    f"{index + 1}/{total_frames}",
                )
                ocr_text = extract_image(frame, report_progress=False).strip()
                if ocr_text:
                    frame_ocr_lines.append(
                        f"### {format_timestamp(index * interval)}\n{ocr_text}"
                    )
        except Exception as exc:  # noqa: BLE001 - frame OCR is useful but optional.
            frame_ocr_lines.append(f"关键帧 OCR 未完成：{exc}")

    parts = [
        "## 视频语音转文字",
        transcript or "未识别到可用语音文字。",
        "",
        "## 视频画面文字 OCR",
        "\n\n".join(frame_ocr_lines) or "未识别到画面文字。",
    ]
    visual_assets = [
        {
            "path": str(frame),
            "kind": "video_frame",
            "label": f"视频时间 {format_timestamp(index * interval)}",
        }
        for index, frame in enumerate(sorted(assets_dir.glob("video_frame_*.jpg")))
    ]
    return "\n".join(parts), visual_assets


def extract_file(path: Path, assets_dir: Path) -> tuple[str, list[dict[str, str]]]:
    """Extract text from one supported file."""

    emit_progress("开始解析文件", 1, path.name)
    extension = path.suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"不支持的文件类型：{extension}")

    if extension in {".md", ".txt", ".csv"}:
        return read_text_file(path), []
    if extension == ".docx":
        return extract_docx(path), extract_docx_visuals(path, assets_dir)
    if extension == ".pdf":
        content = extract_pdf(path)
        return content, render_pdf_pages(path, assets_dir)
    if extension == ".xlsx":
        return extract_xlsx(path), render_xlsx_pages(path, assets_dir)
    if extension == ".pptx":
        content = extract_pptx(path)
        return content, render_ppt_slides(path, assets_dir)
    if extension in IMAGE_EXTENSIONS:
        content = extract_image(path)
        target = save_analysis_image(path, assets_dir / "image_0001.jpg")
        return content, [{"path": str(target), "kind": "image", "label": "原始图片"}]
    if extension in VIDEO_EXTENSIONS:
        return extract_video(path, assets_dir)

    raise ValueError(f"不支持的文件类型：{extension}")


def build_payload(path: Path, assets_dir: Path) -> dict[str, Any]:
    """Build the JSON response consumed by the Node backend."""

    assets_dir.mkdir(parents=True, exist_ok=True)
    content, visual_assets = extract_file(path, assets_dir)
    content = content.strip()
    emit_progress("文件解析完成", 100, path.name)
    return {
        "ok": True,
        "file": str(path),
        "extension": path.suffix.lower(),
        "characters": len(content),
        "content": content,
        "visualAssets": visual_assets,
    }


def parse_args() -> argparse.Namespace:
    """Parse CLI arguments."""

    parser = argparse.ArgumentParser(description="Extract text from a document file.")
    parser.add_argument("file", type=Path)
    parser.add_argument("--assets-dir", type=Path)
    return parser.parse_args()


def main() -> int:
    """Run the extractor and print JSON to stdout."""

    args = parse_args()
    try:
        assets_dir = args.assets_dir or Path(tempfile.mkdtemp(prefix="kb_visual_assets_"))
        payload = build_payload(args.file, assets_dir)
    except Exception as exc:  # noqa: BLE001 - CLI must return structured errors.
        print(
            json.dumps(
                {
                    "ok": False,
                    "file": str(args.file),
                    "error": str(exc),
                },
                ensure_ascii=False,
            ),
            file=sys.stdout,
        )
        return 1

    print(json.dumps(payload, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
